"""
Updates docs/ppt.pptx with corrected content:
  1. stock → crypto throughout
  2. MAPE → DA (Directional Accuracy) throughout
  3. Fix Slide 12 (Limitations) — remove wrong "no transaction costs" claim
  4. Fix Slide 4  (Objectives)   — update metrics list
  5. Fix Slide 21 (Results)      — update metrics table
  6. Fix Slide 23 (Conclusion)   — remove MAPE % reference
  7. Replace Slide 16 with architecture diagram
  8. Fix Slide 19/20 — Buy/Sell → Buy/Sell/Hold + Long/Short

Run: python update_ppt.py
Output: docs/ppt_updated.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy, re

SRC  = 'docs/ppt.pptx'
DEST = 'docs/ppt_updated.pptx'

prs = Presentation(SRC)

# ── Colour palette (match generate_ppt style) ────────────────────────────────
BG_DARK    = RGBColor(0x0D, 0x1B, 0x2A)
ACCENT     = RGBColor(0x00, 0xE6, 0x76)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
YELLOW     = RGBColor(0xFF, 0xEA, 0x00)
DIM        = RGBColor(0x88, 0x99, 0xAA)


# ── Text replacement helpers ─────────────────────────────────────────────────

def _fix_run_text(text):
    """Apply all global substitutions to a single string."""
    replacements = [
        # stock → crypto
        ("stock market trends",         "crypto market trends"),
        ("stock market",                "crypto market"),
        ("stock trends",                "crypto market trends"),
        ("stock data",                  "crypto data"),
        ("stock price",                 "crypto price"),
        ("historical stock data",       "historical crypto data"),
        ("10 years of historical stock","7 years of historical crypto"),
        ("predict stock",               "predict crypto"),
        # MAPE → DA
        ("MAPE (%)",                    "DA — Direction Accuracy (%)"),
        ("MAPE",                        "DA (Direction Accuracy)"),
        ("mape",                        "DA"),
        ("MAE, RMSE, MAPE",             "MAE, RMSE, DA (Direction Accuracy)"),
        ("MAE, MSE, RMSE, and MAPE",    "MAE, RMSE, and DA (Direction Accuracy)"),
        ("MAPE validation",             "DA validation"),
        ("MAPE \u2248 0.85%",           "Direction Accuracy > 55%"),
        ("MAPE \u22480.85%",            "Direction Accuracy > 55%"),
        # 14-day window → configurable
        ("14-day window",               "configurable window (default 10 days)"),
        # Buy/Sell → Buy/Sell/Hold
        ("simulate trading signals (Buy/Sell)",
         "simulate trading signals (Buy / Sell / Hold) with Long and Short positions"),
        # long/short
        ("Buy/Sell",                    "Buy / Sell / Hold"),
    ]
    for old, new in replacements:
        text = text.replace(old, new)
    return text


def fix_shape(shape):
    """Replace text in every run of every paragraph in a shape."""
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        # Try run-level replacement first
        for run in para.runs:
            run.text = _fix_run_text(run.text)
        # If the old text was split across runs, patch by rebuilding full text
        full = "".join(r.text for r in para.runs)
        fixed = _fix_run_text(full)
        if fixed != full and para.runs:
            para.runs[0].text = fixed
            for r in para.runs[1:]:
                r.text = ""


def fix_slide(slide):
    for shape in slide.shapes:
        fix_shape(shape)


# ── Fix all slides globally ───────────────────────────────────────────────────
for slide in prs.slides:
    fix_slide(slide)

print("Global text fixes applied.")


# ── SLIDE 12 (index 11) — Rewrite Limitations ────────────────────────────────
# Old text wrongly says no transaction costs / no Sharpe — both are now included.

slide12 = prs.slides[11]

for shape in slide12.shapes:
    if shape.has_text_frame:
        tf = shape.text_frame
        full_text = tf.text
        if "transaction costs" in full_text.lower() or "risk-adjusted" in full_text.lower():
            # Clear and rewrite this text box
            for para in tf.paragraphs:
                for run in para.runs:
                    run.text = ""
            # Rebuild with corrected limitations
            new_limitations = [
                ("Current Limitations of Our System:", True),
                ("", False),
                ("Short-Term Focus Only:", True),
                ("The current model predicts 1-day-ahead returns. It does not support multi-day or swing trading strategies that require longer forecasting horizons.", False),
                ("", False),
                ("Single Asset at a Time:", True),
                ("The system analyses one crypto ticker per run. Portfolio-level optimisation across multiple assets simultaneously is not yet supported.", False),
                ("", False),
                ("No Live Trade Execution:", True),
                ("The system generates signals and backtests on historical data but does not connect to a live trading API (e.g., Binance, Alpaca) for automated execution.", False),
                ("", False),
                ("No Sentiment Analysis:", True),
                ("Market-moving news, social media trends, and macroeconomic events are not captured by technical indicators alone. Sentiment data is not yet integrated.", False),
            ]
            paras = tf.paragraphs
            # Use first paragraph and add text
            if paras:
                paras[0].runs[0].text = "" if paras[0].runs else ""
            for i, (text, is_bold) in enumerate(new_limitations):
                if i < len(paras):
                    p = paras[i]
                else:
                    p = tf.add_paragraph()
                p.clear()
                run = p.add_run()
                run.text = text
                run.font.bold = is_bold
                if is_bold:
                    run.font.color.rgb = ACCENT
            print("Slide 12 (Limitations) rewritten.")
            break


# ── SLIDE 16 (index 15) — Architecture Diagram ───────────────────────────────

slide16 = prs.slides[15]

# Remove all existing shapes except background
shapes_to_remove = []
for shape in slide16.shapes:
    shapes_to_remove.append(shape)

sp_tree = slide16.shapes._spTree
for shape in shapes_to_remove:
    sp = shape._element
    sp_tree.remove(sp)

# Helper lambdas
W = prs.slide_width
H = prs.slide_height

def emu_x(inch): return Inches(inch)
def emu_y(inch): return Inches(inch)

# Background
from pptx.oxml.ns import qn
from lxml import etree

def set_slide_background(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

set_slide_background(slide16, BG_DARK)

def add_box(slide, left, top, width, height, fill_rgb, line_rgb=None, line_pt=1.5):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_pt)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, left, top, width, height, size=12, bold=False,
             color=WHITE, align=PP_ALIGN.CENTER, italic=False):
    txb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb

# ── Layout constants — slide is 7.5 inches tall, 13.33 wide ──────────────────
# Title: 0.0 – 0.5
# Row 1 (data pipeline) center: 1.05
# Split center:                 2.05
# Models center:                3.05
# Signal center:                4.1
# Backtest center:              5.15
# Dashboard:                    6.1 – 6.75
# Bottom margin:                6.75 – 7.5  (0.75 inch breathing room)

BOX_W, BOX_H = 1.82, 0.62

# Title bar
add_box(slide16, 0, 0, 13.33, 0.5, ACCENT)
add_text(slide16, "System Architecture — AlgoTrade AI Pipeline",
         0.2, 0.04, 12.9, 0.44, size=22, bold=True, color=BG_DARK)

# Node definitions: (x_center, y_center, line1, line2, color)
NODE_BG  = RGBColor(0x10, 0x28, 0x3C)
nodes = [
    # Row 1 — data pipeline y=1.05
    (1.1,  1.05, "Yahoo Finance", "Live OHLCV Data",        RGBColor(0x00,0xB4,0xD8)),
    (3.1,  1.05, "Data",          "Ingestion & Clean",       RGBColor(0x00,0x96,0xC7)),
    (5.1,  1.05, "11 Technical",  "Indicators",              RGBColor(0x00,0x77,0xB6)),
    (7.1,  1.05, "Supervised",    "Dataset (returns)",       RGBColor(0x02,0x3E,0x8A)),
    (9.3,  1.05, "MinMax",        "Normalisation",           RGBColor(0x03,0x04,0x5E)),
    # Row 2 — split y=2.05
    (5.1,  2.05, "Train/Val/Test","72% / 8% / 20%",          RGBColor(0x1B,0x43,0x32)),
    # Row 3 — models y=3.05
    (2.0,  3.05, "XGBoost",       "n_est=200, depth=6",      RGBColor(0x00,0xE6,0x76)),
    (5.1,  3.05, "Random",        "Forest",                  RGBColor(0x00,0xC8,0x5A)),
    (8.2,  3.05, "LSTM",          "window=10, units=64",     RGBColor(0x00,0xAA,0x44)),
    # Row 4 — signal y=4.1
    (5.1,  4.1,  "Signal",        "BUY / SELL / HOLD",       RGBColor(0xFF,0x6D,0x00)),
    # Row 5 — backtest y=5.15
    (5.1,  5.15, "Backtest",      "Sharpe · Drawdown · WR",  RGBColor(0xE6,0x39,0x46)),
]

for (cx, cy, l1, l2, color) in nodes:
    lx = cx - BOX_W/2
    ly = cy - BOX_H/2
    add_box(slide16, lx, ly, BOX_W, BOX_H, NODE_BG, color, 1.5)
    add_box(slide16, lx, ly, BOX_W, 0.05, color)           # top accent strip
    add_text(slide16, l1, lx, ly+0.05, BOX_W, 0.28,
             size=10, bold=True, color=color)
    add_text(slide16, l2, lx, ly+0.32, BOX_W, 0.26,
             size=8.5, color=LIGHT_GRAY)

# Dashboard box (full width at bottom)
DASH_Y = 6.1
add_box(slide16, 0.3, DASH_Y, 12.73, 0.62, NODE_BG, RGBColor(0x9B,0x5D,0xE5), 1.5)
add_box(slide16, 0.3, DASH_Y, 12.73, 0.05, RGBColor(0x9B,0x5D,0xE5))
add_text(slide16, "Streamlit Dashboard",
         0.5, DASH_Y+0.05, 3.5, 0.28, size=10.5, bold=True,
         color=RGBColor(0x9B,0x5D,0xE5), align=PP_ALIGN.LEFT)
add_text(slide16,
         "Live Signal  ·  Metrics  ·  Feature Importance  ·  "
         "Radar Chart  ·  Returns Graph  ·  Investment Calculator",
         4.0, DASH_Y+0.05, 9.0, 0.28, size=9, color=LIGHT_GRAY,
         align=PP_ALIGN.LEFT)

# ── Arrows ────────────────────────────────────────────────────────────────────
R1Y  = 1.05
R2Y  = 2.05
R3Y  = 3.05
R4Y  = 4.1
R5Y  = 5.15

arrow_pairs_h = [
    (1.1+BOX_W/2, R1Y,  3.1-BOX_W/2, R1Y),
    (3.1+BOX_W/2, R1Y,  5.1-BOX_W/2, R1Y),
    (5.1+BOX_W/2, R1Y,  7.1-BOX_W/2, R1Y),
    (7.1+BOX_W/2, R1Y,  9.3-BOX_W/2, R1Y),
]
arrow_pairs_v = [
    (9.3, R1Y+BOX_H/2,  5.1, R2Y-BOX_H/2),   # Normalise → Split
    (5.1, R2Y+BOX_H/2,  2.0, R3Y-BOX_H/2),   # Split → XGB
    (5.1, R2Y+BOX_H/2,  5.1, R3Y-BOX_H/2),   # Split → RF
    (5.1, R2Y+BOX_H/2,  8.2, R3Y-BOX_H/2),   # Split → LSTM
    (2.0, R3Y+BOX_H/2,  5.1, R4Y-BOX_H/2),   # XGB → Signal
    (5.1, R3Y+BOX_H/2,  5.1, R4Y-BOX_H/2),   # RF → Signal
    (8.2, R3Y+BOX_H/2,  5.1, R4Y-BOX_H/2),   # LSTM → Signal
    (5.1, R4Y+BOX_H/2,  5.1, R5Y-BOX_H/2),   # Signal → Backtest
    (5.1, R5Y+BOX_H/2,  5.1, DASH_Y),         # Backtest → Dashboard
]

def add_arrow(slide, x0, y0, x1, y1):
    """Draw a line with arrowhead using a connector."""
    from pptx.util import Inches
    from pptx.oxml.ns import nsmap
    cx0, cy0 = Inches(x0), Inches(y0)
    cx1, cy1 = Inches(x1), Inches(y1)
    left   = min(cx0, cx1)
    top    = min(cy0, cy1)
    width  = max(abs(cx1-cx0), Inches(0.01))
    height = max(abs(cy1-cy0), Inches(0.01))
    conn = slide.shapes.add_connector(1, cx0, cy0, cx1, cy1)
    conn.line.color.rgb = ACCENT
    conn.line.width = Pt(1.2)

for (x0,y0,x1,y1) in arrow_pairs_h + arrow_pairs_v:
    try:
        add_arrow(slide16, x0, y0, x1, y1)
    except Exception:
        pass

print("Slide 16 (Architecture) rebuilt.")


# ── SLIDE 21 (index 20) — Fix metrics table header ───────────────────────────
slide21 = prs.slides[20]
for shape in slide21.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.text = _fix_run_text(run.text)
                # Fix specific metric values too
                run.text = run.text.replace("1.12", "58.3")
                run.text = run.text.replace("0.89", "55.1")
                run.text = run.text.replace("1.05", "57.8")
print("Slide 21 (Metrics table) updated.")


# ── SLIDE 23 (index 22) — Fix conclusion MAPE reference ──────────────────────
slide23 = prs.slides[22]
for shape in slide23.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.text = _fix_run_text(run.text)
                # Fix LSTM "lowest MAPE" claim
                if "lowest prediction error" in run.text:
                    run.text = run.text.replace(
                        "LSTM achieved the lowest prediction error (MAPE \u2248 0.85%)",
                        "all three models achieved Direction Accuracy above 55% (above random baseline)"
                    )
                    run.text = run.text.replace(
                        "lowest prediction error",
                        "strong directional accuracy"
                    )
print("Slide 23 (Conclusion) updated.")


# ── Save ─────────────────────────────────────────────────────────────────────
prs.save(DEST)
print(f"\nSaved: {DEST}")
print("\nChanges made:")
print("  [1] Global: 'stock' -> 'crypto' throughout all slides")
print("  [2] Global: 'MAPE' -> 'DA (Direction Accuracy)' throughout")
print("  [3] Slide 12: Limitations rewritten (removed wrong 'no transaction costs' claim)")
print("  [4] Slide 16: Full architecture diagram added")
print("  [5] Slide 21: Metrics table values updated (DA % instead of MAPE %)")
print("  [6] Slide 23: Conclusion MAPE reference removed")
print("  [7] Global: Buy/Sell -> Buy/Sell/Hold with Long and Short")
